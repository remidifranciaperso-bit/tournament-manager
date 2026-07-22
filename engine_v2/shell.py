"""Coquille PDF Engine V2 (PyMuPDF) — indices slides template, bandeau Live V2."""

from __future__ import annotations

from pathlib import Path

import fitz
from pptx import Presentation

from engine.live_page_map import _slide_text
from engine.pdf_titles import pdf_header_title
from engine_v2.pages.convocations import render_convocations_page
from engine_v2.pages.cover import render_cover_page
from engine_v2.pages.participants import render_participants_page
from engine_v2.pages._layout import (
    draw_page_footer_logo,
    draw_page_header,
    prepare_content_page,
)
from engine_v2.template_registry import export_basename, resolve_template_bundle


def _slide_page_size(cache: dict) -> tuple[float, float]:
    sizes = cache.get("page_sizes") or {}
    if sizes:
        first = next(iter(sizes.values()))
        return float(first["width"]), float(first["height"])
    from engine_v2.config import PAGE_HEIGHT_PT, PAGE_WIDTH_PT

    return PAGE_WIDTH_PT, PAGE_HEIGHT_PT


def _indices_par_role_from_prs(prs: Presentation) -> tuple[set[int], set[int]]:
    participants: set[int] = set()
    convocations: set[int] = set()
    for index, slide in enumerate(prs.slides):
        upper = _slide_text(slide).upper()
        if "CONVOCATION" in upper:
            convocations.add(index)
        elif "PARTICIPANT" in upper:
            participants.add(index)
    return participants, convocations


def _indices_par_role(template_path: Path) -> tuple[set[int], set[int]]:
    prs = Presentation(str(template_path))
    return _indices_par_role_from_prs(prs)


def _titles_by_index(page_map: dict) -> dict[int, str]:
    titles: dict[int, str] = {}
    defaults = {
        "main": "TABLEAU PRINCIPAL",
        "classement": "MATCHS CLASSEMENT",
        "planning": "PLANNING",
        "final": "CLASSEMENT FINAL",
    }
    for section, default in defaults.items():
        for entry in page_map.get(section, []):
            label = pdf_header_title(entry.get("label"), default)
            titles[int(entry["index"])] = label
    return titles


def _render_chrome_shell(
    page: fitz.Page,
    tournoi,
    title: str,
    *,
    base_dir: Path,
    logo_bytes: bytes | None,
    logo_wh: tuple[int, int] | None,
) -> None:
    """Page coquille : bandeau V2 + zone blanche (contenu = capture Live)."""
    prepare_content_page(page)
    draw_page_header(page, tournoi, title, base_dir=base_dir)
    draw_page_footer_logo(
        page,
        logo_bytes=logo_bytes,
        logo_wh=logo_wh,
        club_name=tournoi.club,
        base_dir=base_dir,
        nb_equipes=tournoi.nb_equipes,
    )


def build_v2_composite_shell_pdf(
    *,
    tournoi,
    matchs,
    base_dir: Path,
    logo_bytes: bytes | None = None,
    logo_wh: tuple[int, int] | None = None,
) -> tuple[Path, str]:
    """
    PDF coquille aligné sur le template PPTX (1 page / slide).

    - Garde, participants, convocations : rendu V2 natif (PyMuPDF)
    - Slides tableau/planning/final : bandeau V2 + fond blanc (captures Live)
    """
    base_dir = Path(base_dir)
    template_path, _template_id, cache = resolve_template_bundle(tournoi, base_dir)
    page_map = cache["page_map"]
    width, height = _slide_page_size(cache)
    prs = Presentation(str(template_path))
    participants, convocations = _indices_par_role_from_prs(prs)
    slide_count = len(prs.slides)
    del prs
    titles = _titles_by_index(page_map)

    exports_dir = base_dir / "exports"
    exports_dir.mkdir(parents=True, exist_ok=True)
    pdf_filename = f"{export_basename(tournoi)}-shell.pdf"
    shell_path = exports_dir / pdf_filename

    doc = fitz.open()
    try:
        for slide_index in range(slide_count):
            page = doc.new_page(width=width, height=height)
            if slide_index == 0:
                render_cover_page(
                    page,
                    tournoi,
                    base_dir=base_dir,
                    logo_bytes=logo_bytes,
                    logo_wh=logo_wh,
                )
            elif slide_index in convocations:
                render_convocations_page(
                    page,
                    tournoi,
                    matchs,
                    base_dir=base_dir,
                    logo_bytes=logo_bytes,
                    logo_wh=logo_wh,
                )
            elif slide_index in participants:
                render_participants_page(
                    page,
                    tournoi,
                    base_dir=base_dir,
                    logo_bytes=logo_bytes,
                    logo_wh=logo_wh,
                )
            elif slide_index in titles:
                title = titles[slide_index]
                if title:
                    _render_chrome_shell(
                        page,
                        tournoi,
                        title,
                        base_dir=base_dir,
                        logo_bytes=logo_bytes,
                        logo_wh=logo_wh,
                    )
                else:
                    prepare_content_page(page)
                    draw_page_header(page, tournoi, "", base_dir=base_dir)
                    draw_page_footer_logo(
                        page,
                        logo_bytes=logo_bytes,
                        logo_wh=logo_wh,
                        club_name=tournoi.club,
                        base_dir=base_dir,
                        nb_equipes=tournoi.nb_equipes,
                    )
            else:
                prepare_content_page(page)

        if doc.page_count != slide_count:
            raise RuntimeError("Coquille V2 : nombre de pages incorrect.")
        doc.save(str(shell_path), garbage=4, deflate=True)
    finally:
        doc.close()

    return shell_path, pdf_filename.replace("-shell.pdf", ".pdf")


def _load_logo(logo_path: Path | None) -> tuple[bytes | None, tuple[int, int] | None]:
    if logo_path is None or not logo_path.is_file():
        return None, None
    try:
        from PIL import Image

        raw = logo_path.read_bytes()
        with Image.open(logo_path) as img:
            return raw, (int(img.width), int(img.height))
    except Exception:
        return None, None
