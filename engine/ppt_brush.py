from __future__ import annotations

from io import BytesIO
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation

BRUSH_FONT_NAME = "Grindy Brush"
EMU_PAR_POUCE = 914_400
DPI_RASTER = 200


def _shape_utilise_brush(shape) -> bool:
    if not shape.has_text_frame:
        return False

    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.text.strip() and run.font.name == BRUSH_FONT_NAME:
                return True
    return False


def _couleur_effective(run, shape, slide_height: int) -> tuple[int, int, int]:
    try:
        if run.font.color is not None and run.font.color.rgb is not None:
            rgb = run.font.color.rgb
            return (rgb[0], rgb[1], rgb[2])
    except AttributeError:
        pass

    # Titres dans la bande bleue : texte clair (schéma bg1).
    if shape.top < slide_height * 0.12:
        return (255, 255, 255)

    return (0, 0, 0)


def _taille_pt(shape) -> float:
    taille = 32.0
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.name != BRUSH_FONT_NAME or not run.text.strip():
                continue
            if run.font.size:
                taille = max(taille, run.font.size / 12_700)
    return taille


def _texte_shape(shape) -> str:
    lignes: list[str] = []
    for para in shape.text_frame.paragraphs:
        morceaux = [
            run.text
            for run in para.runs
            if run.font.name == BRUSH_FONT_NAME and run.text
        ]
        if morceaux:
            lignes.append("".join(morceaux).strip())
    return "\n".join(ligne for ligne in lignes if ligne)


def _render_shape_png(shape, font_path: Path, slide_height: int) -> BytesIO | None:
    texte = _texte_shape(shape)
    if not texte:
        return None

    largeur_px = max(1, int(shape.width / EMU_PAR_POUCE * DPI_RASTER))
    hauteur_px = max(1, int(shape.height / EMU_PAR_POUCE * DPI_RASTER))
    taille_pt = _taille_pt(shape)
    taille_px = max(8, int(taille_pt / 72 * DPI_RASTER))

    couleur = (0, 0, 0)
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if run.font.name == BRUSH_FONT_NAME and run.text.strip():
                couleur = _couleur_effective(run, shape, slide_height)
                break

    image = Image.new("RGBA", (largeur_px, hauteur_px), (0, 0, 0, 0))
    draw = ImageDraw.Draw(image)
    police = ImageFont.truetype(str(font_path), taille_px)

    bbox = draw.multiline_textbbox((0, 0), texte, font=police)
    tw = bbox[2] - bbox[0]
    th = bbox[3] - bbox[1]
    x = (largeur_px - tw) / 2 - bbox[0]
    y = (hauteur_px - th) / 2 - bbox[1]
    draw.multiline_text((x, y), texte, font=police, fill=couleur + (255,))

    flux = BytesIO()
    image.save(flux, format="PNG")
    flux.seek(0)
    return flux


def rasteriser_textes_brush(pptx_path: Path, font_path: Path | None = None) -> bool:
    """
    Remplace les zones « Grindy Brush » par des images PNG.

    LibreOffice ne charge pas toujours cette police en headless ; la rasterisation
    garantit le rendu brush identique à l'Engine dans le PDF final.
    """
    pptx_path = Path(pptx_path)
    if font_path is None:
        font_path = pptx_path.resolve().parents[1] / "fonts" / "Grindy Brush.otf"
    font_path = Path(font_path)
    if not font_path.is_file():
        return False

    presentation = Presentation(str(pptx_path))
    slide_height = presentation.slide_height
    modifie = False

    for slide in presentation.slides:
        a_traiter = [shape for shape in slide.shapes if _shape_utilise_brush(shape)]

        for shape in a_traiter:
            png = _render_shape_png(shape, font_path, slide_height)
            if png is None:
                continue

            slide.shapes.add_picture(
                png,
                shape.left,
                shape.top,
                width=shape.width,
                height=shape.height,
            )
            element = shape._element
            element.getparent().remove(element)
            modifie = True

    if modifie:
        presentation.save(str(pptx_path))

    return modifie
