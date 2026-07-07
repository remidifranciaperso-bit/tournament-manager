import re

from engine.ppt_engine import (
    construire_valeurs_globales,
    construire_valeurs_matchs,
    construire_valeurs_planning,
    construire_valeurs_points,
    construire_valeurs_poules,
    valeur_balise_speciale,
)

_STATIC_PREFIXES = ("WIN_", "LOSE_", "SECOND_", "THIRD_")


def _est_balise_statique_masque(nom: str) -> bool:
    return nom.startswith(_STATIC_PREFIXES)


def construire_valeurs_masque_template(template_path) -> dict[str, str]:
    """
    Valeurs pour exporter un masque PNG : libellés WIN_/LOSE_/SECOND_ figés,
    champs dynamiques vidés.
    """
    from pptx import Presentation

    from engine.ppt_engine import parcourir_shapes

    prs = Presentation(str(template_path))
    valeurs: dict[str, str] = {}

    for slide in prs.slides:
        for shape in parcourir_shapes(slide.shapes):
            if getattr(shape, "has_text_frame", False):
                textes = [shape.text_frame.text]
            elif getattr(shape, "has_table", False):
                textes = [
                    cell.text
                    for row in shape.table.rows
                    for cell in row.cells
                ]
            else:
                continue

            for texte in textes:
                for nom in re.findall(r"\{\{([^}]+)\}\}", texte):
                    balise = f"{{{{{nom}}}}}"
                    if balise in valeurs:
                        continue

                    special = valeur_balise_speciale(balise)
                    if special is not None:
                        valeurs[balise] = special
                    elif _est_balise_statique_masque(nom):
                        if nom.startswith("SECOND_"):
                            suffixe = nom.replace("SECOND_", "", 1)
                            valeurs[balise] = f"Deuxième {suffixe}:"
                        elif nom.startswith("THIRD_"):
                            suffixe = nom.replace("THIRD_", "", 1)
                            valeurs[balise] = f"Troisième {suffixe}:"
                        else:
                            valeurs[balise] = ""
                    else:
                        valeurs[balise] = ""

    return valeurs


def construire_champs_live(tournoi, matchs) -> dict[str, str]:
    """Champs dynamiques pour overlay front (clés sans accolades)."""
    valeurs: dict[str, str] = {}
    valeurs.update(construire_valeurs_globales(tournoi))
    valeurs.update(construire_valeurs_matchs(matchs))
    valeurs.update(construire_valeurs_poules(tournoi))
    valeurs.update(construire_valeurs_planning(matchs, tournoi))
    valeurs.update(construire_valeurs_points(tournoi))

    champs: dict[str, str] = {}
    for balise, valeur in valeurs.items():
        cle = balise.strip("{}")
        champs[cle] = str(valeur) if valeur is not None else ""

    return champs
