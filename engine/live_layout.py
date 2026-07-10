import json
import re
from pathlib import Path

from pptx import Presentation

from engine.ppt_engine import parcourir_shapes

SLIDE_W = 9906000
SLIDE_H = 6858000

_TAG = re.compile(r"\{\{([^}]+)\}\}")
_PLANNING_CODE = re.compile(r"^(?:J\d+_)?PL\d+_CODE$")
_DONE_MARKERS = frozenset({"☐", "☑", "✓", "✔", "□"})


def _pct(value: int, total: int) -> float:
    return round(value / total * 100, 3)


def _pct_float(value: float, total: float) -> float:
    return round(value / total * 100, 3)


def _field_dict(
    key: str,
    left: int,
    top: int,
    width: int,
    height: int,
    slide_w: int,
    slide_h: int,
) -> dict:
    return {
        "key": key,
        "left": _pct(left, slide_w),
        "top": _pct(top, slide_h),
        "width": _pct(width, slide_w),
        "height": _pct(height, slide_h),
    }


def _cell_rect(shape, row_idx: int, col_idx: int) -> tuple[int, int, int, int]:
    table = shape.table
    left = shape.left + sum(table.columns[i].width for i in range(col_idx))
    top = shape.top + sum(table.rows[i].height for i in range(row_idx))
    width = table.columns[col_idx].width
    height = table.rows[row_idx].height
    return left, top, width, height


def _append_text_shape_fields(
    fields: list[dict],
    shape,
    slide_w: int,
    slide_h: int,
) -> None:
    if not getattr(shape, "has_text_frame", False):
        return

    text = shape.text_frame.text.strip()
    if not text or "{{" not in text:
        return

    for key in _TAG.findall(text):
        fields.append(
            _field_dict(
                key.strip(),
                shape.left,
                shape.top,
                shape.width,
                shape.height,
                slide_w,
                slide_h,
            )
        )


def _append_table_fields(
    fields: list[dict],
    shape,
    slide_w: int,
    slide_h: int,
) -> None:
    if not getattr(shape, "has_table", False):
        return

    table = shape.table
    nrows = len(table.rows)
    ncols = len(table.columns)
    if nrows == 0 or ncols == 0:
        return

    termine_col = None
    for col_idx in range(ncols):
        header = table.cell(0, col_idx).text.strip().upper()
        if "TERMIN" in header:
            termine_col = col_idx
            break

    for row_idx in range(nrows):
        code_key = None
        done_col = None

        for col_idx in range(ncols):
            text = table.cell(row_idx, col_idx).text.strip()
            tags = _TAG.findall(text)

            for tag in tags:
                key = tag.strip()
                left, top, width, height = _cell_rect(shape, row_idx, col_idx)
                fields.append(
                    _field_dict(key, left, top, width, height, slide_w, slide_h)
                )
                if _PLANNING_CODE.match(key):
                    code_key = key

            if text in _DONE_MARKERS:
                done_col = col_idx

        target_done_col = done_col if done_col is not None else termine_col

        if code_key and target_done_col is not None:
            left, top, width, height = _cell_rect(shape, row_idx, target_done_col)
            fields.append(
                _field_dict(
                    code_key.replace("_CODE", "_DONE"),
                    left,
                    top,
                    width,
                    height,
                    slide_w,
                    slide_h,
                )
            )


def extraire_layout_template(template_path) -> dict[str, list[dict]]:
    prs = Presentation(str(template_path))
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    layout: dict[str, list[dict]] = {}

    for slide_index, slide in enumerate(prs.slides):
        fields: list[dict] = []
        for shape in parcourir_shapes(slide.shapes):
            _append_text_shape_fields(fields, shape, slide_w, slide_h)
            _append_table_fields(fields, shape, slide_w, slide_h)

        if fields:
            layout[str(slide_index)] = fields

    return layout


def extraire_layout_planning(template_path, page_map: dict) -> dict[str, list[dict]]:
    layout = extraire_layout_template(template_path)
    planning_layout: dict[str, list[dict]] = {}

    for entry in page_map.get("planning", []):
        index = str(entry["index"])
        if index in layout:
            planning_layout[index] = layout[index]

    return planning_layout


def _extraire_checkboxes_pdf_page(page) -> list[dict]:
    width = float(page.rect.width)
    height = float(page.rect.height)
    if width <= 0 or height <= 0:
        return []

    boxes: list[dict] = []
    for block in page.get_text("dict")["blocks"]:
        if block.get("type") != 0:
            continue
        for line in block["lines"]:
            for span in line["spans"]:
                if span["text"].strip() not in _DONE_MARKERS:
                    continue
                x0, y0, x1, y1 = span["bbox"]
                boxes.append(
                    {
                        "left": _pct_float(x0, width),
                        "top": _pct_float(y0, height),
                        "width": _pct_float(x1 - x0, width),
                        "height": _pct_float(y1 - y0, height),
                    }
                )

    boxes.sort(key=lambda box: box["top"])
    return boxes


def extraire_checkboxes_pdf(pdf_path, page_indices: list[int]) -> dict[str, list[dict]]:
    import fitz

    doc = fitz.open(str(pdf_path))
    result: dict[str, list[dict]] = {}

    try:
        for index in sorted(set(page_indices)):
            if index < 0 or index >= doc.page_count:
                continue
            boxes = _extraire_checkboxes_pdf_page(doc[index])
            if boxes:
                result[str(index)] = boxes
    finally:
        doc.close()

    return result


def _centrer_hit_area(template_box: dict, pdf_box: dict) -> dict:
    pdf_cx = pdf_box["left"] + pdf_box["width"] / 2
    pdf_cy = pdf_box["top"] + pdf_box["height"] / 2
    width = max(template_box["width"], pdf_box["width"], 3.5)
    height = max(template_box["height"], pdf_box["height"], 3.0)

    return {
        "left": round(pdf_cx - width / 2, 3),
        "top": round(pdf_cy - height / 2, 3),
        "width": round(width, 3),
        "height": round(height, 3),
    }


def calibrer_planning_layout_pdf(
    planning_layout: dict[str, list[dict]],
    pdf_path,
) -> dict[str, list[dict]]:
    """Aligne les zones TERMINÉ sur les cases ☐ du PDF Engine (évite la dérive PPTX)."""
    if not planning_layout:
        return planning_layout

    page_indices = [int(index) for index in planning_layout]
    pdf_boxes = extraire_checkboxes_pdf(pdf_path, page_indices)

    for page_index, slide_fields in planning_layout.items():
        done_fields = sorted(
            (field for field in slide_fields if field["key"].endswith("_DONE")),
            key=lambda field: field["top"],
        )
        page_boxes = pdf_boxes.get(page_index, [])
        if not done_fields or not page_boxes:
            continue

        for field, pdf_box in zip(done_fields, page_boxes):
            hit_area = _centrer_hit_area(field, pdf_box)
            field["left"] = hit_area["left"]
            field["top"] = hit_area["top"]
            field["width"] = hit_area["width"]
            field["height"] = hit_area["height"]

    return planning_layout


def chemin_layout_public(base_dir: Path, template_id: str) -> Path:
    return (
        base_dir
        / "frontend"
        / "public"
        / "live-templates"
        / template_id
        / "layout.json"
    )


def ecrire_layout_public(base_dir: Path, template_id: str, layout: dict) -> Path:
    path = chemin_layout_public(base_dir, template_id)
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(layout, indent=2, ensure_ascii=False), encoding="utf-8")
    return path


def lire_layout_public(base_dir: Path, template_id: str) -> dict[str, list[dict]]:
    path = chemin_layout_public(base_dir, template_id)
    if not path.exists():
        return extraire_layout_template(
            base_dir / "templates bleus" / f"{template_id}.pptx"
        )
    return json.loads(path.read_text(encoding="utf-8"))
