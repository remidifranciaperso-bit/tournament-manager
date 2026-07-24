import re

from pptx import Presentation

from engine.ppt_engine import parcourir_shapes
from engine.live_layout import extraire_layout_template

_CODE_TAG = re.compile(
    r"\{\{([A-Z0-9_]+)_(?:CODE|TERRAIN|HEURE|EQ1|EQ2)\}\}"
)
_JOUR_TAG = re.compile(r"JOUR\s*(\d+)", re.IGNORECASE)
_PLANNING_FRACTION = re.compile(
    r"PLANNING[^0-9]*(\d+\s*/\s*\d+)",
    re.IGNORECASE,
)
_MULTI_DAY_PLANNING_TAG = re.compile(r"\{\{J\d+_PL\d+_")
_LEGACY_PLANNING_TAG = re.compile(r"\{\{PL\d+_")
_LEGACY_PL_CODE = re.compile(r"^PL\d+_CODE$")
_MULTI_DAY_PL_CODE = re.compile(r"^J\d+_PL\d+_CODE$")

_MAIN_PREFIXES = ("P", "H", "Q", "D", "S")
_MAIN_EXACT = frozenset({"F", "PF"})

_FILLED_CLASSEMENT = re.compile(
    r"\bC[0-9]+(?:_[0-9]+)+\b|CLASSEMENT [0-9]"
)
_FILLED_MAIN = re.compile(
    r"(?:^|\n|\s)(?:H[1-8]|Q[1-4]|P[1-8]|D[12]|PF|F|S[0-9]+)(?:\s|\n|$)",
    re.MULTILINE,
)
_FILLED_POULE = re.compile(r"\bPOULE [A-D]\b", re.IGNORECASE)


def _slide_text(slide) -> str:
    parts = []
    for shape in parcourir_shapes(slide.shapes):
        if getattr(shape, "has_text_frame", False):
            text = shape.text_frame.text.strip()
            if text:
                parts.append(text)
    return "\n".join(parts)


def _slide_planning_tag_text(slide) -> str:
    """Texte slide + cellules tableau (balises ``{{PL…}}`` / ``{{J1_PL…}}``)."""
    parts = [_slide_text(slide)]
    for shape in parcourir_shapes(slide.shapes):
        if not getattr(shape, "has_table", False):
            continue
        table = shape.table
        for row in table.rows:
            for cell in row.cells:
                text = cell.text.strip()
                if text:
                    parts.append(text)
    return "\n".join(parts)


def _slide_titles(slide) -> list[str]:
    titres = []
    for shape in parcourir_shapes(slide.shapes):
        if not getattr(shape, "has_text_frame", False):
            continue
        texte = shape.text_frame.text.strip()
        if not texte or "{{" in texte:
            continue
        if len(texte) > 80:
            continue
        titres.append(texte)
    return titres


def _slide_codes(text: str) -> set[str]:
    return set(_CODE_TAG.findall(text))


def _is_main_code(code: str) -> bool:
    if code in _MAIN_EXACT:
        return True
    return any(code.startswith(prefix) for prefix in _MAIN_PREFIXES) and not code.startswith(
        "C"
    )


def _is_classement_code(code: str) -> bool:
    return code.startswith("C")


def _skip_slide(text_upper: str, slide_index: int) -> bool:
    if slide_index == 0:
        return True
    if "PARTICIPANTS" in text_upper:
        return True
    if "CONVOCATIONS" in text_upper:
        return True
    return False


def _classify_content_slide(text: str, text_upper: str) -> str | None:
    if "CLASSEMENT FINAL" in text_upper or (
        "POINTS" in text_upper and "CLASSEMENT" in text_upper
    ):
        return "final"

    if "PLANNING" in text_upper:
        return "planning"

    codes = _slide_codes(text)
    main_codes = [code for code in codes if _is_main_code(code)]
    classement_codes = [code for code in codes if _is_classement_code(code)]

    if main_codes:
        return "main"
    if classement_codes:
        return "classement"

    if _FILLED_CLASSEMENT.search(text):
        return "classement"
    if _FILLED_MAIN.search(text):
        return "main"
    if _FILLED_POULE.search(text):
        return "main"

    return None


def _planning_jour(text_upper: str) -> int | None:
    match = _JOUR_TAG.search(text_upper)
    if match:
        return int(match.group(1))
    return None


def _label_planning(slide) -> str:
    for titre in _slide_titles(slide):
        if "PLANNING" in titre.upper():
            return titre.strip()
    text_upper = _slide_text(slide).upper().replace("–", "-").replace("—", "-")
    match = _PLANNING_FRACTION.search(text_upper)
    if match:
        return f"Planning {match.group(1).replace(' ', '')}"
    return "Planning"


def _label_classement(slide) -> str:
    for titre in _slide_titles(slide):
        upper = titre.upper()
        if "CLASSEMENT" in upper and "FINAL" not in upper:
            return titre
    return "Classement"


_PRELIM_CODE = re.compile(r"^P\d+$")


def _is_prelim_slide(text: str) -> bool:
    codes = _slide_codes(text)
    if not codes:
        return False

    p_codes = {code for code in codes if _PRELIM_CODE.match(code)}
    h_codes = {code for code in codes if re.match(r"^H\d+$", code)}
    q_codes = {code for code in codes if re.match(r"^Q\d+$", code)}

    if p_codes and not h_codes and not q_codes:
        return True

    text_upper = text.upper()
    return any(
        token in text_upper
        for token in ("PRELIMINAIRE", "PRÉLIMINAIRE", "PRELIMIN", "TOUR PRE")
    )


def _label_main_entries(indices: list[int], prs: Presentation) -> list[dict]:
    slides_info = []
    for index in indices:
        slide = prs.slides[index]
        text = _slide_text(slide)
        slides_info.append(
            {
                "index": index,
                "is_prelim": _is_prelim_slide(text),
            }
        )

    bracket = [info for info in slides_info if not info["is_prelim"]]
    bracket_total = len(bracket)
    bracket_pos = 0
    entries = []

    for info in slides_info:
        if info["is_prelim"]:
            label = "Tour préliminaire"
        elif bracket_total == 2:
            label = "Tableau principal" if bracket_pos == 0 else "Partie basse"
            bracket_pos += 1
        elif bracket_total == 1:
            label = "Tableau principal"
        else:
            bracket_pos += 1
            label = f"Partie {bracket_pos}"

        entries.append({"index": info["index"], "label": label})

    return entries


def _label_final(slide) -> str:
    for titre in _slide_titles(slide):
        if "CLASSEMENT FINAL" in titre.upper():
            return "Classement final"
    return "Classement final"


def _as_page_entries(indices: list[int], prs: Presentation, kind: str) -> list[dict]:
    if kind == "main":
        return _label_main_entries(indices, prs)

    entries = []
    for index in indices:
        slide = prs.slides[index]
        if kind == "planning":
            label = _label_planning(slide)
        elif kind == "classement":
            label = _label_classement(slide)
        else:
            label = _label_final(slide)

        entries.append({"index": index, "label": label})

    return entries


def _filter_redundant_planning_entries(
    planning_entries: list[dict],
    prs: Presentation,
) -> list[dict]:
    """
    Retire les slides planning legacy ``PL{n}`` lorsqu'un planning multi-jours
    ``J{jour}_PL{n}`` est présent (sinon doublon Jour 1, ex. « Planning 1/4 »).
    """
    if not planning_entries:
        return planning_entries

    slide_texts = {
        entry["index"]: _slide_planning_tag_text(prs.slides[entry["index"]])
        for entry in planning_entries
    }
    has_multi_day = any(
        _MULTI_DAY_PLANNING_TAG.search(text) for text in slide_texts.values()
    )
    if not has_multi_day:
        return planning_entries

    filtered: list[dict] = []
    for entry in planning_entries:
        text = slide_texts[entry["index"]]
        if _MULTI_DAY_PLANNING_TAG.search(text):
            filtered.append(entry)
            continue
        if _LEGACY_PLANNING_TAG.search(text):
            continue
        filtered.append(entry)
    return filtered


def _filter_redundant_planning_from_layout(
    planning_entries: list[dict],
    planning_layout: dict,
) -> list[dict]:
    """Même filtre que PPTX, à partir des clés ``planning_layout``."""
    if not planning_entries or not planning_layout:
        return planning_entries

    has_multi_day = any(
        _MULTI_DAY_PL_CODE.match(field.get("key", ""))
        for fields in planning_layout.values()
        for field in fields
    )
    if not has_multi_day:
        return planning_entries

    filtered: list[dict] = []
    for entry in planning_entries:
        fields = planning_layout.get(str(entry["index"]), [])
        has_slide_multi = any(
            _MULTI_DAY_PL_CODE.match(field.get("key", "")) for field in fields
        )
        has_slide_legacy = any(
            _LEGACY_PL_CODE.match(field.get("key", "")) for field in fields
        )
        if has_slide_legacy and not has_slide_multi:
            continue
        filtered.append(entry)
    return filtered


def _group_planning_pages_from_layout(
    planning_entries: list[dict],
    planning_layout: dict,
) -> list[list[int]]:
    if not planning_entries:
        return []

    groups: dict[int, list[int]] = {}
    for entry in planning_entries:
        fields = planning_layout.get(str(entry["index"]), [])
        jour: int | None = None
        for field in fields:
            match = re.match(r"^J(\d+)_PL\d+_", field.get("key", ""))
            if match:
                jour = int(match.group(1))
                break
        groups.setdefault(jour or 1, []).append(entry["index"])

    return [groups[jour] for jour in sorted(groups)]


def elaguer_planning_layout(page_map: dict, planning_layout: dict) -> dict:
    """Retire les slides planning exclues du ``page_map``."""
    allowed = {str(entry["index"]) for entry in page_map.get("planning", [])}
    return {
        key: fields for key, fields in planning_layout.items() if key in allowed
    }


def normaliser_page_map_planning(
    page_map: dict,
    *,
    planning_layout: dict | None = None,
    template_path=None,
) -> dict:
    """
    Applique le filtre anti-doublon planning (legacy ``PL{n}`` vs ``J{jour}_PL{n}``).
    Utilisable sur cache figé ou snapshot sans regénérer le PPTX.
    """
    planning = list(page_map.get("planning") or [])
    if not planning:
        return page_map

    if template_path is not None:
        layout = extraire_layout_template(template_path)
        filtered = _filter_redundant_planning_from_layout(planning, layout)
        groups = _group_planning_pages_from_layout(filtered, layout)
    elif planning_layout:
        filtered = _filter_redundant_planning_from_layout(planning, planning_layout)
        groups = _group_planning_pages_from_layout(filtered, planning_layout)
    else:
        return page_map

    if filtered == planning and groups == page_map.get("planning_groups"):
        return page_map

    normalise = dict(page_map)
    normalise["planning"] = filtered
    normalise["planning_groups"] = groups
    return normalise


def _group_planning_pages(
    planning_entries: list[dict],
    prs: Presentation,
) -> list[list[int]]:
    if not planning_entries:
        return []

    entries = []
    for entry in planning_entries:
        text_upper = _slide_text(prs.slides[entry["index"]]).upper()
        entries.append(
            {
                "index": entry["index"],
                "jour": _planning_jour(text_upper),
            }
        )

    if any(entry["jour"] is not None for entry in entries):
        groups: dict[int, list[int]] = {}
        for entry in entries:
            jour = entry["jour"] or 1
            groups.setdefault(jour, []).append(entry["index"])
        return [groups[jour] for jour in sorted(groups)]

    return [[entry["index"] for entry in planning_entries]]


def cartographier_slides_presentation(prs: Presentation) -> dict:
    main: list[int] = []
    classement: list[int] = []
    planning: list[int] = []
    final: list[int] = []

    for index, slide in enumerate(prs.slides):
        text = _slide_text(slide)
        text_upper = text.upper()

        if _skip_slide(text_upper, index):
            continue

        kind = _classify_content_slide(text, text_upper)
        if kind == "main":
            main.append(index)
        elif kind == "classement":
            classement.append(index)
        elif kind == "planning":
            planning.append(index)
        elif kind == "final":
            final.append(index)

    main_entries = _as_page_entries(main, prs, "main")
    classement_entries = _as_page_entries(classement, prs, "classement")
    planning_entries = _filter_redundant_planning_entries(
        _as_page_entries(planning, prs, "planning"),
        prs,
    )
    final_entries = _as_page_entries(final, prs, "final")

    return {
        "main": main_entries,
        "classement": classement_entries,
        "planning": planning_entries,
        "planning_groups": _group_planning_pages(planning_entries, prs),
        "final": final_entries,
    }


def cartographier_slides(pptx_path) -> dict:
    prs = Presentation(str(pptx_path))
    return cartographier_slides_presentation(prs)


def cartographier_pages_live(template_path, pptx_path) -> dict:
    """
    Cartographie fiable : template (balises) puis repli sur le PPTX rempli.
    """
    page_map = cartographier_slides(template_path)

    if page_map["main"] or page_map["classement"]:
        return page_map

    return cartographier_slides(pptx_path)


def indices_depuis_page_map(page_map: dict) -> list[int]:
    indices: set[int] = set()

    for key in ("main", "classement", "planning", "final"):
        for item in page_map.get(key, []):
            if isinstance(item, dict):
                indices.add(int(item["index"]))
            else:
                indices.add(int(item))

    return sorted(indices)
