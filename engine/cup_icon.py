"""Insère une icône coupe PNG sur les libellés vainqueur (LibreOffice ne gère pas le glyphe Unicode)."""

from __future__ import annotations

import re
from io import BytesIO
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Emu

_CUP_PATH = Path(__file__).resolve().parent / "assets" / "cup-winner.png"

_PATTERN_VAINQUEUR = re.compile(
    r"^(?:"
    r"H\d+|Q\d+|D\d+|F|PF|"
    r"Poule [A-D]|"
    r"C\d+_\d+(?:_\d+)?"
    r"):",
    re.IGNORECASE,
)


def _est_libelle_vainqueur(texte: str) -> bool:
    return bool(_PATTERN_VAINQUEUR.match(texte.strip()))


def _remplacer_texte_shape(shape, nouveau: str) -> None:
    tf = shape.text_frame
    if not tf.paragraphs:
        tf.text = nouveau
        return
    first = True
    for para in tf.paragraphs:
        for run in para.runs:
            if first:
                run.text = nouveau
                first = False
            else:
                run.text = ""


def _stream_icone(icon_path: Path) -> BytesIO:
    img = Image.open(icon_path).convert("RGBA")
    buf = BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def inserer_icones_coupe_vainqueur(pptx_path: Path, icon_path: Path | None = None) -> bool:
    icon_path = icon_path or _CUP_PATH
    if not icon_path.is_file():
        return False

    prs = Presentation(str(pptx_path))
    modifie = False
    stream = _stream_icone(icon_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            texte = shape.text_frame.text.strip()
            if not _est_libelle_vainqueur(texte):
                continue

            taille = int(shape.height * 0.82)
            decalage = Emu(28000)
            gauche = max(0, shape.left - taille - decalage)
            haut = shape.top + (shape.height - taille) // 2

            stream.seek(0)
            slide.shapes.add_picture(stream, gauche, haut, width=taille, height=taille)
            modifie = True

    if modifie:
        prs.save(str(pptx_path))

    return modifie
