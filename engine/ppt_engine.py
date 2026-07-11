from copy import deepcopy
from datetime import datetime, timedelta
from pathlib import Path
import re

from engine.points_engine import get_points
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from PIL import Image

# Emojis Noto Color (Dockerfile installe fonts-noto-color-emoji).
ICONE_VAINQUEUR = "🏆 "
ICONE_PERDANT = "❌ "
ICONE_PREMIER = "🏆 "
ICONE_DEUXIEME = "🥈 "
ICONE_TROISIEME = "🥉 "


def format_date(date_str):
    try:
        return datetime.strptime(str(date_str), "%Y-%m-%d").strftime("%d/%m/%y")
    except Exception:
        return str(date_str)


def format_nombre(valeur):
    try:
        nombre = int(valeur)
        return f"{nombre:,}".replace(",", " ")
    except Exception:
        return str(valeur)


def equipe_label_court(equipe):
    if not isinstance(equipe, str):
        return equipe.nom_complet_court()

    texte = equipe.strip()

    if texte.startswith("Vainqueur Poule "):
        return ICONE_PREMIER + texte.replace("Vainqueur ", "") + ":"

    if texte.startswith("Deuxième Poule "):
        return ICONE_DEUXIEME + texte.replace("Deuxième ", "") + ":"

    if texte.startswith("Troisième Poule "):
        return ICONE_TROISIEME + texte.replace("Troisième ", "") + ":"

    if texte.startswith("Second Poule "):
        return ICONE_DEUXIEME + texte.replace("Second ", "") + ":"

    if texte.startswith("Vainqueur "):
        return ICONE_VAINQUEUR + texte.replace("Vainqueur ", "") + ":"

    if texte.startswith("Perdant "):
        return ICONE_PERDANT + texte.replace("Perdant ", "") + ":"

    return texte

def equipe_label_planning(equipe):
    if not isinstance(equipe, str):
        return equipe.nom_court()

    texte = equipe.strip()

    if texte.startswith("Vainqueur Poule "):
        return ICONE_PREMIER + texte.replace("Vainqueur ", "") + ":"

    if texte.startswith("Deuxième Poule "):
        return ICONE_DEUXIEME + texte.replace("Deuxième ", "") + ":"

    if texte.startswith("Troisième Poule "):
        return ICONE_TROISIEME + texte.replace("Troisième ", "") + ":"

    if texte.startswith("Second Poule "):
        return ICONE_DEUXIEME + texte.replace("Second ", "") + ":"

    if texte.startswith("Vainqueur "):
        return ICONE_VAINQUEUR + texte.replace("Vainqueur ", "") + ":"

    if texte.startswith("Perdant "):
        return ICONE_PERDANT + texte.replace("Perdant ", "") + ":"

    return texte


def valeur_balise_speciale(balise):
    nom = balise.strip().replace("{{", "").replace("}}", "")

    if nom.startswith("WIN_"):
        return ICONE_VAINQUEUR + nom.replace("WIN_", "", 1) + ":"

    if nom.startswith("LOSE_"):
        return ICONE_PERDANT + nom.replace("LOSE_", "", 1) + ":"

    if nom.startswith("SECOND_"):
        return ICONE_DEUXIEME + nom.replace("SECOND_", "", 1) + ":"

    if nom.startswith("THIRD_"):
        return ICONE_TROISIEME + nom.replace("THIRD_", "", 1) + ":"

    return None


def remplacer_texte_preserve_style(text_frame, nouveau_texte):
    if not text_frame.paragraphs:
        text_frame.text = str(nouveau_texte)
        return

    p = text_frame.paragraphs[0]

    if not p.runs:
        p.text = str(nouveau_texte)
        return

    p.runs[0].text = str(nouveau_texte)

    for run in p.runs[1:]:
        run.text = ""


def parcourir_shapes(shapes):
    for shape in shapes:
        yield shape

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from parcourir_shapes(shape.shapes)


# Éléments de spPr qui, selon le schéma OOXML, doivent suivre effectLst :
# on insère effectLst juste avant eux pour garder un XML valide.
_APRES_EFFECTLST = (qn("a:scene3d"), qn("a:sp3d"), qn("a:extLst"))


def redresser_connecteurs(prs):
    """
    Redresse les connecteurs obliques en rendu LibreOffice.

    Les connecteurs des templates sont liés à leurs formes via stCxn/endCxn
    (point d'accroche). PowerPoint et LibreOffice ne calculent pas ce point
    d'accroche au même endroit -> LibreOffice dessine une branche en biais.
    On retire la liaison : le connecteur garde sa géométrie stockée (xfrm),
    qui est bien droite et identique au rendu PowerPoint d'origine.
    """

    for slide in prs.slides:
        for shape in parcourir_shapes(slide.shapes):
            element = shape._element

            if element.tag != qn("p:cxnSp"):
                continue

            nv_cxn = element.find(qn("p:nvCxnSpPr"))
            if nv_cxn is None:
                continue

            c_nv_cxn = nv_cxn.find(qn("p:cNvCxnSpPr"))
            if c_nv_cxn is None:
                continue

            for tag in ("a:stCxn", "a:endCxn"):
                for liaison in c_nv_cxn.findall(qn(tag)):
                    c_nv_cxn.remove(liaison)


def supprimer_ombres_theme(prs):
    """
    Neutralise les ombres héritées du thème.

    Les formes des templates n'ont pas d'ombre explicite mais référencent un
    style d'effet du thème (effectRef idx != 0) qui, lui, contient une ombre.
    PowerPoint ne l'affiche pas, mais LibreOffice l'applique -> ombre parasite
    en rendu PDF.

    Deux neutralisations combinées (LibreOffice ignore un effectLst vide et
    applique quand même l'effectRef) :
      1. effectRef idx="0" -> plus aucune référence à un effet du thème ;
      2. <a:effectLst/> vide dans spPr -> formatage local explicite « sans
         effet », qui prime sur la référence de style.
    """

    for slide in prs.slides:
        for shape in parcourir_shapes(slide.shapes):
            element = shape._element

            style = element.find(qn("p:style"))
            if style is not None:
                effect_ref = style.find(qn("a:effectRef"))
                if effect_ref is not None and effect_ref.get("idx") not in (
                    None,
                    "0",
                ):
                    effect_ref.set("idx", "0")

            sp_pr = element.find(qn("p:spPr"))
            if sp_pr is None:
                continue

            if (
                sp_pr.find(qn("a:effectLst")) is not None
                or sp_pr.find(qn("a:effectDag")) is not None
            ):
                continue

            effect_lst = sp_pr.makeelement(qn("a:effectLst"), {})

            insertion = None
            for enfant in sp_pr:
                if enfant.tag in _APRES_EFFECTLST:
                    insertion = enfant
                    break

            if insertion is not None:
                insertion.addprevious(effect_lst)
            else:
                sp_pr.append(effect_lst)


def remplacer_dans_paragraphe(paragraphe, valeurs):
    # On travaille au niveau du paragraphe (et non du text_frame entier) :
    # concaténer tous les paragraphes avec des "\n" puis tout réinjecter dans
    # un seul run insère un saut de ligne DANS le run, ce que LibreOffice rend
    # mal (perte du centrage / du gras). Paragraphe par paragraphe, la
    # structure et l'alignement d'origine sont préservés.
    texte_original = "".join(run.text for run in paragraphe.runs)

    if not texte_original:
        return 0

    nouveau = texte_original
    nb = 0

    for balise, valeur in valeurs.items():
        if balise in nouveau:
            nouveau = nouveau.replace(balise, str(valeur))
            nb += 1

    def repl(match):
        nonlocal nb

        balise = match.group(0)

        special = valeur_balise_speciale(balise)

        if special is not None:
            nb += 1
            return special

        return balise

    nouveau = re.sub(
        r"\{\{(?:WIN|LOSE|SECOND|THIRD)_[A-Z0-9_]+\}\}",
        repl,
        nouveau,
    )

    if nouveau != texte_original:
        if paragraphe.runs:
            paragraphe.runs[0].text = nouveau

            for run in paragraphe.runs[1:]:
                run.text = ""
        else:
            paragraphe.text = nouveau

    return nb


def remplacer_dans_text_frame(text_frame, valeurs):
    if not text_frame.paragraphs:
        return 0

    nb = 0

    for paragraphe in text_frame.paragraphs:
        nb += remplacer_dans_paragraphe(paragraphe, valeurs)

    return nb


def remplacer_balises(prs, valeurs):
    total = 0

    for slide in prs.slides:
        for shape in parcourir_shapes(slide.shapes):

            if getattr(shape, "has_text_frame", False):
                total += remplacer_dans_text_frame(
                    shape.text_frame,
                    valeurs,
                )

            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        total += remplacer_dans_text_frame(
                            cell.text_frame,
                            valeurs,
                        )

    print(f"Balises remplacées : {total}")


def dupliquer_ligne_tableau(
    table,
    ligne_modele_index,
    nombre_lignes_total,
):
    while len(table.rows) < nombre_lignes_total:
        ligne_modele = table._tbl.tr_lst[ligne_modele_index]
        nouvelle_ligne = deepcopy(ligne_modele)
        table._tbl.append(nouvelle_ligne)


def set_cell_text(cell, text):
    remplacer_texte_preserve_style(
        cell.text_frame,
        text,
    )


def slide_contient_texte(slide, texte_recherche):
    for shape in parcourir_shapes(slide.shapes):

        if getattr(shape, "has_text_frame", False):
            if texte_recherche in shape.text_frame.text:
                return True

    return False


def trouver_tables_participants(prs):
    tables = []

    for slide in prs.slides:

        # Préfixe "PARTICIPANT" : couvre "PARTICIPANTS" (H/Mixte) et
        # "PARTICIPANTES" (Femmes), la balise ayant déjà été remplacée.
        if not slide_contient_texte(slide, "PARTICIPANT"):
            continue

        for shape in parcourir_shapes(slide.shapes):

            if getattr(shape, "has_table", False):
                tables.append(shape.table)
                break

    return tables


def trouver_tables_convocations(prs):
    tables = []

    for slide in prs.slides:
        if not slide_contient_texte(slide, "CONVOCATIONS"):
            continue

        for shape in parcourir_shapes(slide.shapes):
            if getattr(shape, "has_table", False):
                tables.append(shape.table)
                break

    return tables

def remplir_table_participants_unique(table, equipes):
    dupliquer_ligne_tableau(
        table=table,
        ligne_modele_index=1,
        nombre_lignes_total=len(equipes) + 1,
    )

    for row_idx, team in enumerate(equipes, start=1):
        # Format template participants classique :
        # NOM J1 | CLASSEMENT J1 | NOM J2 | CLASSEMENT J2 | POIDS | TS

        set_cell_text(
            table.cell(row_idx, 0),
            team.joueur1,
        )

        set_cell_text(
            table.cell(row_idx, 1),
            format_nombre(team.classement_j1),
        )

        set_cell_text(
            table.cell(row_idx, 2),
            team.joueur2,
        )

        set_cell_text(
            table.cell(row_idx, 3),
            format_nombre(team.classement_j2),
        )

        set_cell_text(
            table.cell(row_idx, 4),
            format_nombre(team.poids),
        )

        set_cell_text(
            table.cell(row_idx, 5),
            team.ts_label(),
        )


_MOTIF_JOUEUR = re.compile(r"joueur", re.IGNORECASE)


def _feminiser_joueur(texte):
    def _remplacer(match):
        mot = match.group(0)
        if mot.isupper():
            return "JOUEUSE"
        if mot[:1].isupper():
            return "Joueuse"
        return "joueuse"

    return _MOTIF_JOUEUR.sub(_remplacer, texte)


def detecter_genre(tournoi):
    genre = getattr(tournoi, "genre_tournoi", None)

    if genre is None:
        type_base = str(getattr(tournoi, "type_tournoi", "")).strip()
        if " - " in type_base:
            morceaux = type_base.split(" - ")
            if len(morceaux) > 1:
                genre = morceaux[1].strip()

    return genre


def feminiser_entetes_participants(prs):
    """
    Sur la page Participants, remplace « Joueur 1/2 » par « Joueuse 1/2 »
    (en respectant la casse d'origine) pour les tournois féminins.
    Agit uniquement sur la ligne d'en-tête du tableau des participants.
    """
    for table in trouver_tables_participants(prs):
        if not table.rows:
            continue

        for cell in table.rows[0].cells:
            for paragraphe in cell.text_frame.paragraphs:
                for run in paragraphe.runs:
                    if _MOTIF_JOUEUR.search(run.text):
                        run.text = _feminiser_joueur(run.text)


def remplir_table_participants(prs, tournoi):
    tables = trouver_tables_participants(prs)

    if not tables:
        return

    if detecter_genre(tournoi) == "Femmes":
        feminiser_entetes_participants(prs)

    equipes = sorted(
        tournoi.equipes,
        key=lambda e: e.ts,
    )

    # 32 équipes : participants sur 2 pages
    if len(tables) >= 2 and len(equipes) > 24:
        remplir_table_participants_unique(
            tables[0],
            equipes[:16],
        )

        remplir_table_participants_unique(
            tables[1],
            equipes[16:],
        )

    else:
        remplir_table_participants_unique(
            tables[0],
            equipes,
        )


def construire_valeurs_globales(tournoi):
    type_base = str(tournoi.type_tournoi).strip()
    genre = getattr(tournoi, "genre_tournoi", None)

    if " - " in type_base:
        morceaux = type_base.split(" - ")
        type_base = morceaux[0].strip()

        if genre is None and len(morceaux) > 1:
            genre = morceaux[1].strip()

    if genre:
        type_affichage = f"{type_base} {genre}"
    else:
        type_affichage = type_base

    nb_terrains = len(tournoi.terrains)
    texte_terrains = (
        f"{nb_terrains} TERRAIN"
        if nb_terrains == 1
        else f"{nb_terrains} TERRAINS"
    )

    nb_equipes = tournoi.nb_equipes
    texte_equipes = (
        f"{nb_equipes} ÉQUIPE"
        if nb_equipes == 1
        else f"{nb_equipes} ÉQUIPES"
    )

    if genre == "Femmes":
        participants = "PARTICIPANTES"
    else:
        participants = "PARTICIPANTS"

    return {
        "{{TYPE}}": type_affichage,
        "{{DATE}}": format_date(tournoi.date_tournoi),
        "{{HEURE}}": tournoi.heure_debut,
        "{{NB_TERRAINS}}": texte_terrains,
        "{{NB_EQUIPES}}": texte_equipes,
        "{{PARTICIPANTS}}": participants,
        "{{CLUB}}": tournoi.club,
    }

def construire_valeurs_matchs(matchs):
    valeurs = {}

    for match in matchs:
        code = match.code

        valeurs[f"{{{{{code}_CODE}}}}"] = match.code
        valeurs[f"{{{{{code}_TERRAIN}}}}"] = match.terrain or ""
        valeurs[f"{{{{{code}_HEURE}}}}"] = match.heure or ""
        valeurs[f"{{{{{code}_EQ1}}}}"] = equipe_label_court(match.equipe1)
        valeurs[f"{{{{{code}_EQ2}}}}"] = equipe_label_court(match.equipe2)

    return valeurs


def construire_valeurs_poules(tournoi):
    valeurs = {}

    exemptes = getattr(tournoi, "exemptes", [])
    poules = getattr(tournoi, "poules", {})

    for i, equipe in enumerate(exemptes, start=1):
        valeurs[f"{{{{EXEMPT_{i}_EQ}}}}"] = equipe_label_court(equipe)

    for nom_poule, equipes in poules.items():
        for i, equipe in enumerate(equipes, start=1):
            valeurs[f"{{{{POULE_{nom_poule}_{i}_EQ}}}}"] = equipe_label_court(equipe)

    valeurs["{{WIN_POULE_A}}"] = ICONE_PREMIER + "Poule A:"
    valeurs["{{WIN_POULE_B}}"] = ICONE_PREMIER + "Poule B:"
    valeurs["{{WIN_POULE_C}}"] = ICONE_PREMIER + "Poule C:"
    valeurs["{{WIN_POULE_D}}"] = ICONE_PREMIER + "Poule D:"

    valeurs["{{SECOND_POULE_A}}"] = ICONE_DEUXIEME + "Poule A:"
    valeurs["{{SECOND_POULE_B}}"] = ICONE_DEUXIEME + "Poule B:"
    valeurs["{{SECOND_POULE_C}}"] = ICONE_DEUXIEME + "Poule C:"
    valeurs["{{SECOND_POULE_D}}"] = ICONE_DEUXIEME + "Poule D:"

    return valeurs
def construire_convocations(
    matchs,
    date_tournoi,
    heure_debut_tournoi="00:00",
    heures_debut_jours=None,
):

    convocations = {}

    tournoi_multi_jours = (
        len(
            set(
                getattr(m, "jour", 1)
                for m in matchs
            )
        ) > 1
    )

    try:
        date_base = datetime.strptime(str(date_tournoi), "%Y-%m-%d")
    except Exception:
        date_base = None

    try:
        heure_limite = datetime.strptime(
            str(heure_debut_tournoi),
            "%H:%M",
        ).time()
    except Exception:
        heure_limite = datetime.strptime("00:00", "%H:%M").time()

    if not heures_debut_jours:
        heures_debut_jours = [heure_debut_tournoi]

    def heure_debut_du_jour(jour):
        # jour est 1-based ; on récupère l'heure de début propre à ce jour.
        index = max(0, jour - 1)

        if index < len(heures_debut_jours):
            valeur = heures_debut_jours[index]
        else:
            valeur = heures_debut_jours[-1]

        try:
            return datetime.strptime(str(valeur), "%H:%M").time()
        except Exception:
            return heure_limite

    def infos_match(match):
        jour = getattr(match, "jour", 1)

        try:
            heure_match = datetime.strptime(match.heure, "%H:%M").time()
        except Exception:
            heure_match = heure_debut_du_jour(jour)

        decalage_jour = jour - 1

        # Bascule "après minuit" : un match dont l'heure est antérieure au
        # début de SON jour appartient à la nuit suivante (ex. finale à 00:40
        # après un début à 18:00). On compare au début du jour concerné, et
        # non au début global du tournoi : sinon un jour 2 qui démarre plus
        # tôt que le jour 1 (ex. 09:00 < 18:00) serait décalé à tort de +1.
        if heure_match < heure_debut_du_jour(jour):
            decalage_jour += 1

        return decalage_jour, heure_match

    matchs_tries = sorted(
        matchs,
        key=lambda m: (
            infos_match(m)[0],
            infos_match(m)[1],
            m.terrain or "",
            getattr(m, "ordre_planning", 9999),
        ),
    )

    for match in matchs_tries:

        decalage_jour, heure_match = infos_match(match)

        if date_base:
            date_match = (
                date_base + timedelta(days=decalage_jour)
            ).strftime("%d/%m/%y")
        else:
            date_match = f"J{decalage_jour + 1}"

        if tournoi_multi_jours:
            texte_heure = f"{match.heure} ({date_match})"
        else:
            texte_heure = match.heure

        for equipe in [match.equipe1, match.equipe2]:

            if isinstance(equipe, str):
                continue

            nom = equipe.nom_complet()

            if nom not in convocations:
                convocations[nom] = (
                    decalage_jour,
                    heure_match,
                    texte_heure,
                )

    return [
        (nom, texte)
        for nom, (_, _, texte) in sorted(
            convocations.items(),
            key=lambda x: (
                x[1][0],
                x[1][1],
            ),
        )
    ]
    
def remplir_table_convocations_unique(table, convocations):
    dupliquer_ligne_tableau(
        table=table,
        ligne_modele_index=1,
        nombre_lignes_total=len(convocations) + 1,
    )

    for row_idx, (nom, heure) in enumerate(convocations, start=1):
        set_cell_text(table.cell(row_idx, 0), nom)
        set_cell_text(table.cell(row_idx, 1), heure)


def remplir_table_convocations(prs, tournoi, matchs):
    tables = trouver_tables_convocations(prs)

    if not tables:
        return

    convocations = construire_convocations(
        matchs,
        tournoi.date_tournoi,
        tournoi.heure_debut,
        getattr(tournoi, "heures_debut_jours", None),
    )

    # 32 équipes : convocations sur 2 pages
    if len(tables) >= 2 and len(convocations) > 24:
        remplir_table_convocations_unique(
            tables[0],
            convocations[:16],
        )

        remplir_table_convocations_unique(
            tables[1],
            convocations[16:],
        )

    else:
        remplir_table_convocations_unique(
            tables[0],
            convocations,
        )


def ajouter_valeurs_planning(valeurs, prefix, match):

    valeurs[f"{{{{{prefix}_CODE}}}}"] = match.code
    valeurs[f"{{{{{prefix}_HEURE}}}}"] = match.heure or ""
    valeurs[f"{{{{{prefix}_TERRAIN}}}}"] = match.terrain or ""
    valeurs[f"{{{{{prefix}_TOUR}}}}"] = match.tour

    valeurs[f"{{{{{prefix}_EQ1}}}}"] = (
        equipe_label_planning(match.equipe1)
    )

    valeurs[f"{{{{{prefix}_EQ2}}}}"] = (
        equipe_label_planning(match.equipe2)
    )


def construire_valeurs_planning(matchs, tournoi=None):

    valeurs = {}

    matchs_tries = sorted(
        matchs,
        key=lambda m: (
            getattr(m, "ordre_planning", None)
            if getattr(m, "ordre_planning", None) is not None
            else m.ordre
        ),
    )

    for i, match in enumerate(matchs_tries, start=1):
        ajouter_valeurs_planning(
            valeurs,
            f"PL{i}",
            match,
        )

    if tournoi is None:
        return valeurs

    nb_equipes = tournoi.nb_equipes
    nb_jours = tournoi.nb_jours
    mode_tournoi = tournoi.mode_tournoi

    groupes = []

    # ==========================
    # ELIMINATION DIRECTE
    # ==========================

    if mode_tournoi == "Élimination directe" and nb_equipes == 20:
        if nb_jours == 2:
            groupes = [
                matchs_tries[0:22],
                matchs_tries[22:40],
            ]

        elif nb_jours == 3:
            groupes = [
                matchs_tries[0:14],
                matchs_tries[14:28],
                matchs_tries[28:40],
            ]

    elif mode_tournoi == "Élimination directe" and nb_equipes == 24:
        if nb_jours == 2:
            groupes = [
                matchs_tries[0:28],
                matchs_tries[28:52],
            ]

        elif nb_jours == 3:
            groupes = [
                matchs_tries[0:20],
                matchs_tries[20:38],
                matchs_tries[38:52],
            ]

    elif mode_tournoi == "Élimination directe" and nb_equipes == 32:
        if nb_jours == 2:
            groupes = [
                matchs_tries[0:48],
                matchs_tries[48:80],
            ]

        elif nb_jours == 3:
            groupes = [
                matchs_tries[0:32],
                matchs_tries[32:52],
                matchs_tries[52:80],
            ]

    # ==========================
    # POULES + TABLEAU FINAL
    # ==========================

    elif mode_tournoi == "Poules + tableau final" and nb_equipes == 20:
        if nb_jours == 2:
            groupes = [
                matchs_tries[0:24],
                matchs_tries[24:36],
            ]

        elif nb_jours == 3:
            groupes = [
                matchs_tries[0:12],
                matchs_tries[12:24],
                matchs_tries[24:36],
            ]

    elif mode_tournoi == "Poules + tableau final" and nb_equipes == 24:
        if nb_jours == 2:
            groupes = [
                matchs_tries[0:24],
                matchs_tries[24:56],
            ]

        elif nb_jours == 3:
            groupes = [
                matchs_tries[0:12],
                matchs_tries[12:24],
                matchs_tries[24:56],
            ]

    elif mode_tournoi == "Poules + tableau final" and nb_equipes == 32:
        if nb_jours == 2:
            groupes = [
                matchs_tries[0:40],
                matchs_tries[40:72],
            ]

        elif nb_jours == 3:
            groupes = [
                matchs_tries[0:20],
                matchs_tries[20:40],
                matchs_tries[40:72],
            ]

    for numero_jour, groupe in enumerate(groupes, start=1):
        for numero_match, match in enumerate(groupe, start=1):
            ajouter_valeurs_planning(
                valeurs,
                f"J{numero_jour}_PL{numero_match}",
                match,
            )

    return valeurs

def construire_valeurs_points(tournoi):
    valeurs = {}

    points = get_points(
        tournoi.type_tournoi,
        tournoi.nb_equipes,
    )

    for place in range(1, tournoi.nb_equipes + 1):
        valeurs[f"{{{{PTS{place}}}}}"] = format_nombre(
            points.get(place, "")
        )

    return valeurs

def balises_restantes(prs):
    restantes = []

    for slide_num, slide in enumerate(prs.slides, start=1):

        for shape in parcourir_shapes(slide.shapes):

            if getattr(shape, "has_text_frame", False):
                texte = shape.text_frame.text

                if "{{" in texte and "}}" in texte:
                    restantes.append(
                        (slide_num, texte.strip())
                    )

            if getattr(shape, "has_table", False):

                for row in shape.table.rows:
                    for cell in row.cells:

                        texte = cell.text

                        if "{{" in texte and "}}" in texte:
                            restantes.append(
                                (slide_num, texte.strip())
                            )

    return restantes

def remplir_template(
    template_path,
    output_path,
    tournoi,
    matchs,
    logo_path=None,
):
    prs = Presentation(template_path)

    remplacer_logo(
        prs=prs,
        logo_path=logo_path,
        club=tournoi.club,
    )

    valeurs = {}

    valeurs.update(
        construire_valeurs_globales(tournoi)
    )

    valeurs.update(
        construire_valeurs_matchs(matchs)
    )

    valeurs.update(
        construire_valeurs_poules(tournoi)
    )

    valeurs.update(
        construire_valeurs_planning(
            matchs,
            tournoi,
        )
    )

    valeurs.update(
        construire_valeurs_points(
            tournoi,
        )
    )

    remplacer_balises(
        prs,
        valeurs,
    )

    remplir_table_participants(
        prs,
        tournoi,
    )

    remplir_table_convocations(
        prs,
        tournoi,
        matchs,
    )

    supprimer_ombres_theme(prs)

    redresser_connecteurs(prs)

    restantes = balises_restantes(prs)

    if restantes:
        print()
        print("===== BALISES RESTANTES =====")

        for slide_num, texte in restantes:
            print(
                f"Diapo {slide_num} : {texte}"
            )

    prs.save(output_path)
    del prs, valeurs

def remplir_template_8(
    template_path,
    output_path,
    tournoi,
    matchs,
):
    remplir_template(
        template_path=template_path,
        output_path=output_path,
        tournoi=tournoi,
        matchs=matchs,
    )

def remplacer_logo(prs, logo_path=None, club=""):
    for slide in prs.slides:
        for shape in list(parcourir_shapes(slide.shapes)):

            if not getattr(shape, "has_text_frame", False):
                continue

            if "{{LOGO}}" not in shape.text_frame.text:
                continue

            box_left = shape.left
            box_top = shape.top
            box_w = shape.width
            box_h = shape.height

            if logo_path:
                marge = 0.05

                with Image.open(logo_path) as img:
                    img_w, img_h = img.size

                max_w = int(box_w * (1 - marge * 2))
                max_h = int(box_h * (1 - marge * 2))

                ratio = min(
                    max_w / img_w,
                    max_h / img_h,
                )

                logo_w = int(img_w * ratio)
                logo_h = int(img_h * ratio)

                logo_left = int(box_left + (box_w - logo_w) / 2)
                logo_top = int(box_top + (box_h - logo_h) / 2)

                element = shape._element
                element.getparent().remove(element)

                slide.shapes.add_picture(
                    str(logo_path),
                    logo_left,
                    logo_top,
                    width=logo_w,
                    height=logo_h,
                )

            else:
                remplacer_texte_preserve_style(
                    shape.text_frame,
                    club,
                )

                # Sans logo, la zone de texte du club ne doit pas afficher
                # sa bordure (le trait w="0" du template se rend comme un
                # rectangle visible en LibreOffice). On la neutralise en
                # gardant le texte du club organisateur.
                shape.line.fill.background()

def analyser_template(fichier):
    prs = Presentation(fichier)

    print()
    print("==============================")
    print("BALISES TROUVÉES")
    print("==============================")
    print()

    for i, slide in enumerate(prs.slides, start=1):

        print(f"===== DIAPO {i} =====")

        for shape in parcourir_shapes(slide.shapes):

            if getattr(shape, "has_text_frame", False):

                texte = shape.text.strip()

                if "{{" in texte and "}}" in texte:
                    print(texte)

            if getattr(shape, "has_table", False):

                for row in shape.table.rows:
                    for cell in row.cells:

                        texte = cell.text.strip()

                        if "{{" in texte and "}}" in texte:
                            print(texte)

        print()