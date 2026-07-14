"""Tampon logo PDF — après LibreOffice (évite l'OOM PPTX).

Le logo est inséré aux emplacements du gabarit ``{{LOGO}}`` du template
(grand en tête sur la page de garde, petit en pied sur les pages de contenu),
en respectant le ratio d'aspect (contain, centré) — jamais étiré.
"""

from __future__ import annotations

import gc
import shutil
import uuid
from pathlib import Path

import fitz

ENGINE_FOOTER_RATIO = 0.042
ENGINE_FOOTER_LOGO_WIDTH_RATIO = 1 / 3
PDF_LOGO_MAX_PX = 220

# Emplacement garde par défaut (templates sans {{LOGO}}, ex. 32 éq.) :
# grand, en haut, centré — fractions (x0, y0, x1, y1), calqué sur les autres templates.
GARDE_LOGO_BOX_DEFAUT = (0.27, 0.055, 0.72, 0.165)

# Type d'une boîte {{LOGO}} en fractions de la diapo : (x0, y0, x1, y1).
LogoBox = tuple[float, float, float, float]


def _zone_logo_pied(page: fitz.Page) -> fitz.Rect:
    """Bande pied historique (repli si le template n'expose pas de {{LOGO}})."""
    rect = page.rect
    footer_top = rect.height * (1 - ENGINE_FOOTER_RATIO)
    side = rect.width * (1 - ENGINE_FOOTER_LOGO_WIDTH_RATIO) / 2
    return fitz.Rect(side, footer_top, rect.width - side, rect.height)


def contain_rect(zone: fitz.Rect, img_w: float, img_h: float) -> fitz.Rect:
    """Sous-rectangle de ``zone`` respectant le ratio du logo (contain, centré)."""
    if img_w <= 0 or img_h <= 0 or zone.width <= 0 or zone.height <= 0:
        return zone
    scale = min(zone.width / img_w, zone.height / img_h)
    draw_w = img_w * scale
    draw_h = img_h * scale
    x0 = zone.x0 + (zone.width - draw_w) / 2
    y0 = zone.y0 + (zone.height - draw_h) / 2
    return fitz.Rect(x0, y0, x0 + draw_w, y0 + draw_h)


def logo_boxes_par_diapo(template_path: Path) -> dict[int, list[LogoBox]]:
    """Emplacements ``{{LOGO}}`` du template, par index de diapo (fractions)."""
    try:
        from pptx import Presentation
    except Exception:
        return {}

    try:
        prs = Presentation(str(template_path))
    except Exception:
        return {}

    slide_w = float(prs.slide_width or 0)
    slide_h = float(prs.slide_height or 0)
    if slide_w <= 0 or slide_h <= 0:
        return {}

    def parcourir(shapes):
        for shape in shapes:
            yield shape
            if getattr(shape, "shape_type", None) == 6:  # groupe
                yield from parcourir(shape.shapes)

    boxes: dict[int, list[LogoBox]] = {}
    for index, slide in enumerate(prs.slides):
        rects: list[LogoBox] = []
        for shape in parcourir(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            if "{{LOGO}}" not in (shape.text_frame.text or ""):
                continue
            rects.append(
                (
                    shape.left / slide_w,
                    shape.top / slide_h,
                    (shape.left + shape.width) / slide_w,
                    (shape.top + shape.height) / slide_h,
                )
            )
        if rects:
            boxes[index] = rects
    return boxes


def _remplacer_fichier_atomique(cible: Path, source: Path) -> None:
    """Remplace cible par source (même volume requis pour rename atomique)."""
    cible = Path(cible)
    source = Path(source)
    try:
        source.replace(cible)
    except OSError:
        shutil.copy2(source, cible)
        source.unlink(missing_ok=True)


def _dimensions_logo(logo_path: Path) -> tuple[int, int]:
    try:
        from PIL import Image

        with Image.open(logo_path) as img:
            return int(img.width), int(img.height)
    except Exception:
        return 1, 1


def appliquer_logo_sur_pdf(
    pdf_path: Path,
    logo_path: Path,
    logo_boxes: dict[int, list[LogoBox]] | None = None,
) -> None:
    """Insère le logo aux emplacements ``{{LOGO}}`` du template (contain).

    ``logo_boxes`` : emplacements par index de diapo (voir ``logo_boxes_par_diapo``).
    Sans ``logo_boxes`` non vide, repli sur la bande pied historique (contain).
    """
    from engine.logo_prepare import preparer_logo_fichier

    pdf_path = Path(pdf_path).resolve()
    logo_path = Path(logo_path)
    if not pdf_path.is_file() or not logo_path.is_file():
        return

    preparer_logo_fichier(logo_path, max_px=PDF_LOGO_MAX_PX)
    logo_bytes = logo_path.read_bytes()
    if len(logo_bytes) < 32:
        return

    img_w, img_h = _dimensions_logo(logo_path)

    # Fichier temporaire dans le même dossier que le PDF (évite errno 18 sur Render).
    temp_path = pdf_path.parent / f".{pdf_path.stem}-logo-{uuid.uuid4().hex}.pdf"

    mode_defaut = not logo_boxes

    doc = fitz.open(str(pdf_path))
    try:
        for index, page in enumerate(doc):
            rect = page.rect
            if not mode_defaut:
                # Gabarit {{LOGO}} du template : overlay transparent (les zones
                # ont été vidées lors du remplissage, rien à masquer).
                cibles = [
                    (
                        fitz.Rect(
                            rect.x0 + fx0 * rect.width,
                            rect.y0 + fy0 * rect.height,
                            rect.x0 + fx1 * rect.width,
                            rect.y0 + fy1 * rect.height,
                        ),
                        False,
                    )
                    for (fx0, fy0, fx1, fy1) in logo_boxes.get(index, [])
                ]
            elif index == 0:
                # Template sans {{LOGO}} (ex. 32 éq.) : logo grand en tête de garde.
                fx0, fy0, fx1, fy1 = GARDE_LOGO_BOX_DEFAUT
                cibles = [
                    (
                        fitz.Rect(
                            rect.x0 + fx0 * rect.width,
                            rect.y0 + fy0 * rect.height,
                            rect.x0 + fx1 * rect.width,
                            rect.y0 + fy1 * rect.height,
                        ),
                        False,
                    )
                ]
            else:
                # Pages de contenu : bande pied sur fond blanc (masque le libellé).
                cibles = [(_zone_logo_pied(page), True)]

            for zone, effacer_fond in cibles:
                if effacer_fond:
                    page.draw_rect(zone, color=None, fill=(1, 1, 1), overlay=True)
                dest = contain_rect(zone, img_w, img_h)
                page.insert_image(dest, stream=logo_bytes, keep_proportion=True)

        doc.save(str(temp_path), garbage=4, deflate=True)
    finally:
        doc.close()
        gc.collect()

    try:
        _remplacer_fichier_atomique(pdf_path, temp_path)
    finally:
        temp_path.unlink(missing_ok=True)
