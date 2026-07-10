"""Test de remplissage des templates 24 équipes (templates bleus).

Vérifie que le moteur remplit correctement chaque template (avec et sans poules,
1J/2J/3J) sans laisser de placeholder {{...}} non substitué.
Ne fait pas la conversion PDF (lente) : patch de convertir_pptx_en_pdf.
"""
import re
import sys
from pathlib import Path

BASE = Path(__file__).resolve().parent
sys.path.insert(0, str(BASE))

import engine.tournament_engine as te
from pptx import Presentation

# On court-circuite la conversion PDF pour ne tester que le remplissage pptx.
te.convertir_pptx_en_pdf = lambda pptx_path, output_dir: pptx_path

EXCEL = BASE / "excels test" / "test24equipes_48participants_complet.xlsx"

CONFIGS = [
    ("Élimination directe", 1),
    ("Élimination directe", 2),
    ("Élimination directe", 3),
    ("Poules + tableau final", 1),
    ("Poules + tableau final", 2),
    ("Poules + tableau final", 3),
]

PLACEHOLDER_RE = re.compile(r"\{\{[^}]+\}\}")


def placeholders_restants(pptx_path):
    prs = Presentation(str(pptx_path))
    restants = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        for m in PLACEHOLDER_RE.findall(run.text or ""):
                            restants.add(m)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for m in PLACEHOLDER_RE.findall(cell.text or ""):
                            restants.add(m)
    return restants


def main():
    print(f"Excel de test : {EXCEL.name}")
    print(f"Existe : {EXCEL.exists()}\n")
    ok = 0
    for mode, jours in CONFIGS:
        label = f"24 éq. · {mode} · {jours}J"
        heures = ["09:00", "09:00", "09:00"][:jours]
        try:
            pdf, _snapshot = te.generate_tournament(
                excel_path=str(EXCEL),
                club="Club Test",
                date_tournoi="2026-09-12",
                type_tournoi="P500 - Messieurs",
                heure_debut="09:00",
                duree_match=40,
                terrains=["Terrain 1", "Terrain 2", "Terrain 3", "Terrain 4"],
                terrain_principal="Terrain 1",
                base_dir=str(BASE),
                mode_tournoi=mode,
                nb_jours=jours,
                heures_debut_jours=heures,
                genre_tournoi="Messieurs",
            )
            restants = placeholders_restants(pdf)
            if restants:
                print(f"[⚠]  {label}")
                print(f"      pptx : {Path(pdf).name}")
                print(f"      placeholders NON remplis ({len(restants)}) :")
                for p in sorted(restants):
                    print(f"        - {p}")
            else:
                print(f"[OK] {label}  →  {Path(pdf).name}")
                ok += 1
        except Exception as e:
            print(f"[ERREUR] {label}")
            print(f"      {type(e).__name__}: {e}")
    print(f"\n{ok}/{len(CONFIGS)} variantes remplies sans placeholder résiduel.")


if __name__ == "__main__":
    main()
